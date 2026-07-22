from pptx import Presentation
prs = Presentation('/home/parajuli/repos/2025_learning_beamforming_rsma_code/docs/EE_presentation.pptx')
print(f'Slide count: {len(prs.slides._sldIdLst)}')
for i, slide in enumerate(prs.slides):
    titles = [s.text_frame.text.strip().splitlines()[0] for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    pics = sum(1 for s in slide.shapes if s.shape_type == 13)
    print(i + 1, titles[0] if titles else '(no title)', f'[{pics} picture(s)]')
