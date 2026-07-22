from pptx import Presentation
from pptx.util import Inches, Pt

FIG_PATH = '/home/parajuli/repos/2025_learning_beamforming_rsma_code/src/energy_efficiency/datellite.png'
EQ_DIR = '/home/parajuli/repos/2025_learning_beamforming_rsma_code/docs/presentation_assets'
OUT_PATH = '/home/parajuli/repos/2025_learning_beamforming_rsma_code/docs/EE_presentation.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]
TITLE_SLIDE = prs.slide_layouts[0]


def add_title(slide, text, top=Inches(0.35), size=30):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    return box


def add_bullets(slide, bullets, top, width=Inches(12.0), font_size=19, left=Inches(0.7)):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(font_size - level * 2)
        p.space_after = Pt(9)
    return box


def add_equation(slide, name, top, left=Inches(1.15), height=Inches(0.62)):
    path = f'{EQ_DIR}/{name}.png'
    slide.shapes.add_picture(path, left, top, height=height)
    return height


def add_table(slide, rows, top, left=Inches(0.6), width=Inches(12.1), col_widths=None, font_size=13, header_font_size=14):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, Inches(0.32 * len(rows)))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(header_font_size if r == 0 else font_size)
                p.font.bold = (r == 0)
    return table_shape


# =================================================================
# Slide 1: Title
# =================================================================
slide = prs.slides.add_slide(TITLE_SLIDE)
slide.shapes.title.text = "Energy-Efficient Robust Beamforming for Satellite Downlink"
subtitle = slide.placeholders[1]
subtitle.text = "Formulating a new maximization problem on our satellite downlink system"

# =================================================================
# Slide 2: Motivation -- why energy efficiency
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Motivation: Why Energy Efficiency?")
add_bullets(slide, [
    ("Real satellites have limited onboard power (solar/battery) and amplifier heat budgets", 0),
    ("Maximizing raw data rate alone always pushes transmit power to its full budget, "
     "regardless of the marginal rate that extra power actually buys", 0),
    ("The real question is different:", 0),
    ("not “how much rate can we get,” but “how much rate do we get per watt actually spent”", 1),
    ("This calls for a new optimization problem — built on the same satellite system "
     "our earlier work already established", 0),
], top=Inches(1.5), font_size=20)

# =================================================================
# Slide 3: Motivation -- why a learning-based approach
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Motivation: Why a Learning-Based Approach?")
add_bullets(slide, [
    ("Classic (“conventional”) precoders — e.g. MMSE, robust SLNR — work by solving a math "
     "formula by hand, derived specifically for one objective and one error assumption", 0),
    ("Our new objective (rate per watt) combined with realistic channel-knowledge errors makes "
     "that hand-derived math extremely hard — in general, provably intractable", 0),
    ("Instead of deriving a new formula from scratch, we let an agent learn a good strategy "
     "directly from repeated trial and error — a Reinforcement Learning (RL) approach", 0),
    ("This is also flexible: the same learning approach adapts automatically to different antenna "
     "counts, user counts, or error levels — no new math derivation needed each time", 0),
], top=Inches(1.5), font_size=19)

# =================================================================
# Slide 4: What is SAC?
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "What Is Soft Actor-Critic (SAC)?")
add_bullets(slide, [
    ("Reinforcement Learning in one line: an agent tries something, sees a “reward” for how well "
     "it worked, and gradually adjusts to do better — automated trial and error", 0),
    ("Two neural networks work together:", 0),
    ("Actor — proposes a precoding strategy (what signal weights to transmit); the decision-maker", 1),
    ("Critic — judges how good that strategy turned out to be; a coach scoring each attempt, "
     "guiding the Actor's next try", 1),
    ("“Soft” = the agent is deliberately kept a little unpredictable/exploratory, so it keeps "
     "searching for a better strategy instead of settling early on a mediocre one", 0),
    ("Why it fits us: our precoding weights are continuous-valued, not a small fixed menu of "
     "choices — and the Actor/Critic never need to know the math form of the objective in "
     "advance, only a reward number after each attempt", 0),
], top=Inches(1.35), font_size=18)

# =================================================================
# Slide 5: System model (the shared baseline system)
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "System Model")
add_bullets(slide, [
    ("Single LEO satellite, Uniform Linear Array (ULA), N = 16 antennas", 0),
    ("K = 3 single-antenna ground users", 0),
    ("Satellite altitude d₀ = 600 km   |   Sat. gain G_Sat = 20 dBi   |   User gain G_Usr = 0 dBi", 0),
    ("Carrier wavelength λ = 15 cm (f = 2 GHz)   |   Inter-antenna spacing dₙ = 3λ/2 ≈ 22.5 cm", 0),
    ("Noise power σₙ² ≈ 6×10⁻¹³ W   |   Transmit power budget P = 100 W", 0),
    ("Same LoS channel model and CSIT error model (AoD error + large-scale fading) established "
     "in our earlier work — the baseline system this new problem is built on", 0),
], top=Inches(1.3), font_size=16)
slide.shapes.add_picture(FIG_PATH, Inches(4.9), Inches(3.15), height=Inches(3.85))

# =================================================================
# Slide 6: Baseline objective (established in earlier work)
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Baseline: Maximize Sum Rate")
add_bullets(slide, [
    ("Our earlier work on this system established the following objective:", 0),
], top=Inches(1.4), font_size=19)
add_equation(slide, 'rbar', top=Inches(2.0), left=Inches(1.15), height=Inches(0.5))
add_bullets(slide, [
    ("where the instantaneous sum rate over all K users is:", 0),
], top=Inches(2.65), font_size=19)
add_equation(slide, 'sumrate', top=Inches(3.25), left=Inches(1.15), height=Inches(0.85))
add_bullets(slide, [
    ("subject to a total radiated power budget P", 0),
    ("The precoder (MMSE, robust SLNR, or our earlier learned SAC precoder) rescales its "
     "output to land exactly on this power budget", 0),
    ("Under this rate-only objective, using more transmit power is never worse — so the "
     "optimum always saturates the full power budget", 0),
    ("This is the baseline we now build a new problem on top of", 0),
], top=Inches(4.35), font_size=17)

# =================================================================
# Slide 7: Our new objective function
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Our New Objective: Energy Efficiency")
add_bullets(slide, [
    ("We add an electrical power consumption model on top of the baseline system:", 0),
], top=Inches(1.4), font_size=18)
add_equation(slide, 'ptotal', top=Inches(2.0), left=Inches(1.15), height=Inches(0.58))
add_bullets(slide, [
    ("our values: η_PA = 0.60 (amplifier efficiency), P_c = 1.0 W/antenna × N=16 → 16 W static draw", 1),
    ("and formulate a new maximization problem — Energy Efficiency:", 0),
], top=Inches(2.85), font_size=18)
add_equation(slide, 'ee', top=Inches(3.85), left=Inches(1.15), height=Inches(0.62))
add_bullets(slide, [
    ("same power budget P = 100 W as the baseline", 1),
    ("Precoder now only rescales DOWN when over budget (“clip-only”), instead of always "
     "normalizing exactly to the boundary — this lets the policy choose to use less than "
     "the full budget when that is more efficient", 0),
], top=Inches(4.7), font_size=18)

# =================================================================
# Slide 8: How we solve it -- SAC in practice
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "How We Solve It: SAC in Practice", size=28)
add_bullets(slide, [
    ("Each simulation step, in a loop:", 0),
    ("channel estimate  →  Actor network  →  proposed precoding matrix W  →  reward  →  "
     "store & learn", 1),
    ("Rate-over-power is a ratio, awkward for RL to optimize directly — Dinkelbach's transform "
     "turns it into a subtraction, folded directly into the reward:", 0),
], top=Inches(1.35), font_size=17)
add_equation(slide, 'dinkelbach', top=Inches(2.85), left=Inches(1.15), height=Inches(0.42))
add_bullets(slide, [
    ("λ updated online, as a running average over a 5,000-step window", 1),
    ("Actor and Critic: 4 fully-connected layers of 512 units each, with batch normalization "
     "between layers to keep training stable", 0),
    ("Trained with batch size 1024, drawing from a replay buffer of 100,000 past experiences, "
     "Adam optimizer", 0),
    ("Trained separately at 3 CSIT error levels: Δε = 0.0, 0.025, 0.05 — "
     "13,000 episodes × 1,000 steps ≈ 13M simulation steps per checkpoint", 0),
], top=Inches(3.55), font_size=17)

# =================================================================
# Slide 9: Results -- power comparison + robustness
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Results: Power Savings and Robustness")
add_bullets(slide, [
    ("MMSE and the regular (100 W) SAC checkpoint always spend the full power budget; our "
     "Energy-Efficient checkpoints do not — and stay robust as CSIT error grows:", 0),
], top=Inches(1.15), font_size=17)

slide.shapes.add_picture(f'{EQ_DIR}/power_comparison_new.png', Inches(0.25), Inches(1.75), height=Inches(4.35))
slide.shapes.add_picture(f'{EQ_DIR}/rate_panel_final.png', Inches(7.5), Inches(1.75), height=Inches(4.35))

cap1 = slide.shapes.add_textbox(Inches(0.25), Inches(6.15), Inches(6.9), Inches(1.1))
cap1.text_frame.word_wrap = True
cap1.text_frame.text = ("MMSE and regular SAC always transmit at the full 100 W budget. Our "
                         "Energy-Efficient SAC checkpoints use only ~56-59 W (41-44% saved).")
cap1.text_frame.paragraphs[0].font.size = Pt(14)

cap2 = slide.shapes.add_textbox(Inches(7.5), Inches(6.15), Inches(5.5), Inches(1.1))
cap2.text_frame.word_wrap = True
cap2.text_frame.text = ("The checkpoint trained at higher CSIT error degrades more gracefully and "
                         "overtakes the others as error grows — clean crossovers around error ≈ 0.03–0.07.")
cap2.text_frame.paragraphs[0].font.size = Pt(14)

# =================================================================
# Slide 10: Results -- spatial behavior (beampattern)
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Results: What This Looks Like Spatially")
add_bullets(slide, [
    ("Comparing where the regular (100 W) SAC checkpoint and the Energy-Efficient checkpoint "
     "actually point their antenna beams, for the same users:", 0),
], top=Inches(1.15), font_size=17)
slide.shapes.add_picture(f'{EQ_DIR}/beampattern21_final.png', Inches(2.9), Inches(1.9), height=Inches(4.3))
cap3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.9))
cap3.text_frame.word_wrap = True
cap3.text_frame.text = ("Both checkpoints aim at the same 3 users — the Energy-Efficient checkpoint's "
                         "beams are just sharper and taller, delivering the same coverage with less "
                         "spread-out (and therefore less wasted) transmit energy.")
cap3.text_frame.paragraphs[0].font.size = Pt(14)

# =================================================================
# Slide 11: Appendix -- full parameter reference (paper vs. our code)
# =================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Appendix: Parameter Reference (Paper vs. Our Code)", size=26)

rows = [
    ("Parameter", "Paper (Table I)", "Our code (config.py)"),
    ("Noise power σₙ²", "6×10⁻¹³ W", "6.02×10⁻¹³ W"),
    ("Transmit power P", "100 W", "100 W"),
    ("Satellite altitude d₀", "600 km", "600 km"),
    ("Antenna Nr. N", "{10, 16}", "16"),
    ("User Nr. K", "3", "3"),
    ("Overall Sat. Gain G_Sat", "20 dBi", "20 dBi"),
    ("Wavelength λ", "15 cm", "15 cm (f = 2 GHz)"),
    ("Gain per user G_Usr", "0 dBi", "0 dBi"),
    ("Inter-antenna distance dₙ", "3λ/2", "3λ/2 ≈ 22.5 cm"),
    ("SGD optimizer", "Adam", "Adam"),
    ("Training batch size B", "1024", "1024"),
    ("Init. LR CrNN, AcNN", "1e-4, 1e-5", "8.8e-6, 4.2e-5 (Optuna-tuned)"),
    ("Learning buffer size", "100,000", "100,000"),
    ("Inference / Learning ratio", "10 : 1", "10 : 1"),
    ("L2 scales α_Q̂, α_μ", "0.1", "0.01"),
    ("log Entropy scale α_e", "Variable", "Adaptive (target entropy fixed at 1.0)"),
]
add_table(slide, rows, top=Inches(1.15), font_size=13, header_font_size=14,
          col_widths=[Inches(4.3), Inches(3.5), Inches(4.3)])

label = slide.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(3.0), Inches(0.5))
label.text_frame.text = "Free-space path loss (paper Eq. 2):"
label.text_frame.paragraphs[0].font.size = Pt(12)
label.text_frame.paragraphs[0].font.italic = True
add_equation(slide, 'fspl', top=Inches(6.62), left=Inches(3.5), height=Inches(0.5))

note2 = slide.shapes.add_textbox(Inches(6.4), Inches(6.78), Inches(6.3), Inches(0.7))
tf2 = note2.text_frame
tf2.word_wrap = True
tf2.text = ("EE-only additions, no analogue in the paper: η_PA=0.60, P_c=1.0 W/antenna (16 W total), "
            "Dinkelbach λ window=5000 steps, Δε ∈ {0.0, 0.025, 0.05}, discount γ=0.0 (bandit-style)")
tf2.paragraphs[0].font.size = Pt(11)
tf2.paragraphs[0].font.italic = True

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
