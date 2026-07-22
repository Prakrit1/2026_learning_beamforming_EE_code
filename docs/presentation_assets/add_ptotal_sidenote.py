"""
Adds a small side-note text box next to the P_total(W) equation on the
"Our new objective: Energy Efficiency" slide, explaining eta_PA and P_c.

Note on terminology (flagged to the user, not silently changed): eta_PA is
the amplifier EFFICIENCY (a fraction in (0,1], e.g. 0.60), not
"inefficiency" -- matches EE_formulation.tex eq. (3) and every other mention
of it in this project (config.py, build_ppt7.py's own slide text). And P_c
is the PER-ANTENNA circuit power (the total circuit draw is N*P_c) -- the
P_total(W) equation's underbrace already labels N*P_c as "circuit draw" as a
whole, with P_c itself being the per-antenna figure. Both corrected here.
"""
import datetime
import shutil

from pptx import Presentation
from pptx.util import Emu, Pt

REPO = '/home/parajuli/repos/2025_learning_beamforming_rsma_code'
DOCS = f'{REPO}/docs'
PPTX_PATH = f'{DOCS}/Energy_Efficiency.pptx'

backup_path = f'{DOCS}/Energy_Efficiency_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx'
shutil.copy2(PPTX_PATH, backup_path)
print(f'Backup saved: {backup_path}')

prs = Presentation(PPTX_PATH)

target_slide = None
for slide in prs.slides:
    for shp in slide.shapes:
        if shp.has_text_frame and 'Our new objective: Energy Efficiency' in shp.text_frame.text:
            target_slide = slide
            break
    if target_slide is not None:
        break

if target_slide is None:
    raise RuntimeError('Could not find the "Our new objective: Energy Efficiency" slide')

pics = [s for s in target_slide.shapes if s.shape_type == 13]
pic_ptotal = next(s for s in pics if s.name == 'Picture 1')  # P_total(W) equation, unchanged

eq_right = pic_ptotal.left + pic_ptotal.width
eq_vcenter = pic_ptotal.top + pic_ptotal.height // 2

box_left = eq_right + Emu(228600)  # ~0.25in gap
box_width = Emu(3350000)
box_height = Emu(700000)
box_top = eq_vcenter - box_height // 2

box = target_slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
tf = box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = 'η_PA = amplifier efficiency'
p1.font.size = Pt(14)
p1.font.italic = True

p2 = tf.add_paragraph()
p2.text = 'P_c = per-antenna circuit power'
p2.font.size = Pt(14)
p2.font.italic = True

prs.save(PPTX_PATH)
print(f'Saved: {PPTX_PATH}')
print(f'Side-note box placed at left={box_left}, top={box_top}, width={box_width}, height={box_height}')
