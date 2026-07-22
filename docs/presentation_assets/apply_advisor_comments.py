"""
Applies Alea's three review comments left as floating text boxes in
docs/Energy_Efficiency.pptx, cross-checked against docs/EE_formulation.tex
(the formal restatement of the reference paper's baseline problem):

1. Slide "System model": italicize math variables (N, K, d0, f, P).
2. Slide "Sum rate Maximization" + the other title-only divider slides
   (originally slides 4, 5, 7, 9, 11): deleted, per "these inbetween slides
   are not necessary for that kind of a presentation".
3. Slide "Baseline: Sum Rate Maximization":
   - sum-rate equation: Gamma_k replaced with the explicit SINR fraction
     (matches EE_formulation.tex eq. 1 / reference paper eq. 6-7).
   - objective equation: rewritten as a proper maximization problem,
     max_W Rbar = E{R} s.t. ||W||_F^2 <= P (matches EE_formulation.tex eq. 2).
   All three comment text boxes are removed once addressed.
"""
import copy
import datetime
import re
import shutil

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

REPO = '/home/parajuli/repos/2025_learning_beamforming_rsma_code'
DOCS = f'{REPO}/docs'
ASSETS = f'{DOCS}/presentation_assets'
PPTX_PATH = f'{DOCS}/Energy_Efficiency.pptx'

# ---------------------------------------------------------------------------
# 1. Render the two replacement equations, matching this project's existing
#    equation-image style (matplotlib mathtext, fontset='cm', 400dpi,
#    transparent bg, tight bbox -- see rbar.png / image11.png metadata).
#    NOTE: text.usetex is deliberately left False, consistent with every
#    other plotting script in src/energy_efficiency/ (dvipng isn't installed
#    on this machine, so usetex=True is broken here; mathtext's 'cm' fontset
#    reproduces the same Computer Modern look without needing a LaTeX call).
# ---------------------------------------------------------------------------
plt.rcParams['mathtext.fontset'] = 'cm'
plt.rcParams['text.usetex'] = False


def render_eq(tex, outpath, fontsize=38):
    fig = plt.figure()
    fig.text(0, 0, tex, fontsize=fontsize)
    fig.savefig(outpath, dpi=400, transparent=True, bbox_inches='tight', pad_inches=0.05)
    plt.close(fig)


EQ_OBJECTIVE = (
    r'$\max_{\mathbf{W}} \;\; \bar{R} = \mathbb{E}\{R\} '
    r'\quad \mathrm{s.t.} \quad \|\mathbf{W}\|_F^2 \leq P$'
)
EQ_SUMRATE_SINR = (
    r'$R = \sum_{k=1}^{K} \log_2(1 + \mathrm{SINR}_k)$'
)

objective_png = f'{ASSETS}/rbar_max_new.png'
sumrate_png = f'{ASSETS}/sumrate_sinr_new.png'
render_eq(EQ_OBJECTIVE, objective_png, fontsize=38)
render_eq(EQ_SUMRATE_SINR, sumrate_png, fontsize=34)
print(f'Rendered {objective_png}')
print(f'Rendered {sumrate_png}')

# ---------------------------------------------------------------------------
# 2. Restore from the pristine pre-edit backup (if this script already ran
#    once, PPTX_PATH is no longer the original 14-slide deck), then take a
#    fresh timestamped backup before editing again.
# ---------------------------------------------------------------------------
PRISTINE_PATH = f'{DOCS}/Energy_Efficiency_pristine.pptx'
FIRST_RUN_BACKUP = f'{DOCS}/Energy_Efficiency_backup_20260717_101421.pptx'
import os
if not os.path.exists(PRISTINE_PATH):
    # PPTX_PATH itself was already edited by a prior run of this script; the
    # true pre-edit original is the very first run's timestamped backup.
    shutil.copy2(FIRST_RUN_BACKUP, PRISTINE_PATH)
    print(f'Seeded pristine copy from {FIRST_RUN_BACKUP}')
shutil.copy2(PRISTINE_PATH, PPTX_PATH)
print(f'Restored {PPTX_PATH} from pristine copy before re-editing')

backup_path = f'{DOCS}/Energy_Efficiency_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx'
shutil.copy2(PPTX_PATH, backup_path)
print(f'Backup saved: {backup_path}')

prs = Presentation(PPTX_PATH)


def find_shape(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def italicize_tokens_in_paragraph(paragraph, tokens):
    pattern = re.compile(
        r'(?<![A-Za-z0-9])(' + '|'.join(re.escape(t) for t in tokens) + r')(?![A-Za-z0-9])'
    )
    for run in list(paragraph.runs):
        text = run.text
        if not pattern.search(text):
            continue
        r_elem = run._r
        parent = r_elem.getparent()
        idx = list(parent).index(r_elem)

        segments = []
        last = 0
        for m in pattern.finditer(text):
            if m.start() > last:
                segments.append((text[last:m.start()], False))
            segments.append((m.group(0), True))
            last = m.end()
        if last < len(text):
            segments.append((text[last:], False))

        new_elems = []
        for seg_text, is_var in segments:
            new_r = copy.deepcopy(r_elem)
            t_elem = new_r.find(qn('a:t'))
            t_elem.text = seg_text
            if is_var:
                rPr = new_r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = new_r.makeelement(qn('a:rPr'), {})
                    new_r.insert(0, rPr)
                rPr.set('i', '1')
            new_elems.append(new_r)

        for i, ne in enumerate(new_elems):
            parent.insert(idx + i, ne)
        parent.remove(r_elem)


def replace_picture(shape, slide_part, new_img_path, new_height_emu, new_top_emu=None):
    w_px, h_px = Image.open(new_img_path).size
    aspect = w_px / h_px
    new_cx = int(new_height_emu * aspect)
    old_left = shape.left
    old_top = shape.top if new_top_emu is None else new_top_emu

    image_part, rId = slide_part.get_or_add_image_part(new_img_path)
    blip = shape._element.blipFill.find(qn('a:blip'))
    blip.set(qn('r:embed'), rId)

    src_rect = shape._element.blipFill.find(qn('a:srcRect'))
    if src_rect is not None:
        shape._element.blipFill.remove(src_rect)

    shape.left = Emu(old_left)
    shape.top = Emu(old_top)
    shape.width = Emu(new_cx)
    shape.height = Emu(new_height_emu)


# ---------------------------------------------------------------------------
# 3. Fix content BEFORE deleting any slides (indices below are the ORIGINAL,
#    pre-deletion 1-based slide numbers, captured as objects so later
#    deletions of other slides can't invalidate these references).
# ---------------------------------------------------------------------------
slide_system_model = prs.slides[2]   # slide 3: "System model"
slide_baseline = prs.slides[5]       # slide 6: "Baseline: Sum Rate Maximization"

# --- Slide 3: italicize math variables, remove the comment box ---
body = find_shape(slide_system_model, 'Textplatzhalter 5')
tokens = ['N', 'K', 'd₀', 'f', 'P']
for para in body.text_frame.paragraphs:
    italicize_tokens_in_paragraph(para, tokens)

comment_3 = find_shape(slide_system_model, 'TextBox 2')
if comment_3 is not None:
    remove_shape(comment_3)

# --- Slide 6: swap in the two new equations, remove both comment boxes ---
pic_objective = find_shape(slide_baseline, 'Picture 1')  # rbar.png -> objective
pic_sumrate = find_shape(slide_baseline, 'Picture 2')    # sumrate.png -> sum rate

replace_picture(pic_objective, slide_baseline.part, objective_png,
                 new_height_emu=600075, new_top_emu=1590699)
replace_picture(pic_sumrate, slide_baseline.part, sumrate_png,
                 new_height_emu=750000, new_top_emu=2584581)

for name in ('TextBox 8', 'TextBox 9'):
    box = find_shape(slide_baseline, name)
    if box is not None:
        remove_shape(box)

print('Content fixes applied to slides 3 and 6.')

# ---------------------------------------------------------------------------
# 4. Delete the title-only "inbetween" divider slides: originally slides
#    4, 5, 7, 9, 11 (1-based) -> 0-based indices 3, 4, 6, 8, 10.
#    Must delete highest index first so earlier indices stay valid.
# ---------------------------------------------------------------------------
def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


for idx in sorted([3, 4, 6, 8, 10], reverse=True):
    delete_slide(prs, idx)

print(f'Deleted 5 divider slides. Remaining slide count: {len(prs.slides._sldIdLst)}')

prs.save(PPTX_PATH)
print(f'Saved: {PPTX_PATH}')
