from pptx import Presentation

PPTX_PATH = '/home/parajuli/repos/2025_learning_beamforming_rsma_code/docs/Energy_Efficiency.pptx'
prs = Presentation(PPTX_PATH)
print(f'Slide count: {len(prs.slides._sldIdLst)}')
for i, slide in enumerate(prs.slides):
    print('=' * 15, 'SLIDE', i + 1, '=' * 15)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f'  [{shape.name}] {shape.text_frame.text!r}')
        elif shape.shape_type == 13:
            print(f'  [{shape.name}] PICTURE left={shape.left} top={shape.top} w={shape.width} h={shape.height}')

print()
print('--- slide 3 body runs (checking italics) ---')
slide3 = prs.slides[2]
for shp in slide3.shapes:
    if shp.name == 'Textplatzhalter 5':
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                print(f'  italic={r.font.italic!s:5} text={r.text!r}')
