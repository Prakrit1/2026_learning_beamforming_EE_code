"""
Follow-up to apply_advisor_comments.py. On the "Our new objective: Energy
Efficiency" slide, the EE ratio was shown as a bare definition,
EE(W) = R(W)/P_total(W), with no constraint. Rewrite it as a proper
maximization problem, matching the treatment already given to the baseline
slide's objective (max_W Rbar = E{R} s.t. ||W||_F^2 <= P), per
EE_formulation.tex eq. (4):

    max_W  EE(W) = R(W)/P_total(W)  s.t.  ||W||_F^2 <= P

Same constraint form as the baseline -- the difference (per user instruction)
is behavioral, not notational: the EE precoder is clip-only (eq. 5 in
EE_formulation.tex) and generically does NOT saturate this constraint,
unlike the baseline's unconditional rescale-to-boundary. That distinction is
already stated in this slide's existing bullet text ("Precoder now only
rescales DOWN when over budget..."), so it is not re-encoded in the equation
itself -- the equation's job is just to state the problem correctly.
"""
import datetime
import shutil

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

REPO = '/home/parajuli/repos/2025_learning_beamforming_rsma_code'
DOCS = f'{REPO}/docs'
ASSETS = f'{DOCS}/presentation_assets'
PPTX_PATH = f'{DOCS}/Energy_Efficiency.pptx'

plt.rcParams['mathtext.fontset'] = 'cm'
plt.rcParams['text.usetex'] = False


def render_eq(tex, outpath, fontsize=38):
    fig = plt.figure()
    fig.text(0, 0, tex, fontsize=fontsize)
    fig.savefig(outpath, dpi=400, transparent=True, bbox_inches='tight', pad_inches=0.05)
    plt.close(fig)


EQ_EE_MAX = (
    r'$\max_{\mathbf{W}} \;\; \mathrm{EE}(\mathbf{W}) = '
    r'\frac{R(\mathbf{W})}{P_{\mathrm{total}}(\mathbf{W})} '
    r'\quad \mathrm{s.t.} \quad \|\mathbf{W}\|_F^2 \leq P$'
)
ee_max_png = f'{ASSETS}/ee_max_new.png'
render_eq(EQ_EE_MAX, ee_max_png, fontsize=38)
print(f'Rendered {ee_max_png}')

backup_path = f'{DOCS}/Energy_Efficiency_backup_{datetime.datetime.now():%Y%m%d_%H%M%S}.pptx'
shutil.copy2(PPTX_PATH, backup_path)
print(f'Backup saved: {backup_path}')

prs = Presentation(PPTX_PATH)


def replace_picture(shape, slide_part, new_img_path, new_height_emu, new_top_emu=None):
    w_px, h_px = Image.open(new_img_path).size
    aspect = w_px / h_px
    new_cx = int(new_height_emu * aspect)
    old_left = shape.left
    old_top = shape.top if new_top_emu is None else new_top_emu

    image_part, rId = slide_part.get_or_add_image_part(new_img_path)
    blip = shape._element.blipFill.find(qn('a:blip'))
    blip.set(qn('r:embed'), rId)

    src_rect = shape._element.blipFill.find(qn('a:srcRect'))
    if src_rect is not None:
        shape._element.blipFill.remove(src_rect)

    shape.left = Emu(old_left)
    shape.top = Emu(old_top)
    shape.width = Emu(new_cx)
    shape.height = Emu(new_height_emu)


target_slide = None
for slide in prs.slides:
    for shp in slide.shapes:
        if shp.has_text_frame and 'Our new objective: Energy Efficiency' in shp.text_frame.text:
            target_slide = slide
            break
    if target_slide is not None:
        break

if target_slide is None:
    raise RuntimeError('Could not find the "Our new objective: Energy Efficiency" slide')

pics = [s for s in target_slide.shapes if s.shape_type == 13]
# Picture 1 = P_total(W) (unchanged), Picture 2 = EE(W) ratio (to replace) --
# identify by name to be explicit rather than relying on shape order.
pic_ee = next(s for s in pics if s.name == 'Picture 2')

old_h = pic_ee.height
old_top = pic_ee.top
growth = 142875  # same delta used for the baseline slide's max_W objective
new_h = old_h + growth
new_top = old_top - growth // 2

replace_picture(pic_ee, target_slide.part, ee_max_png, new_height_emu=new_h, new_top_emu=new_top)
print(f'Replaced EE objective equation: new height={new_h}, new top={new_top}')

prs.save(PPTX_PATH)
print(f'Saved: {PPTX_PATH}')
